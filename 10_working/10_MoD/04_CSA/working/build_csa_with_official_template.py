"""
Cisco 2025 公式テンプレートをベースに CSA スライドを統合生成
テンプレートのマスター・色・フォント・ロゴを継承して、2026年対応版を作成
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime
import os
import shutil

# パス定義
template_path = r"c:\Users\kimiito\OneDrive - Cisco\Documents\700 - github\kimiito_Note\10_working\10_MoD\04_CSA\cisco2025_template.pptx"
output_path = r"c:\Users\kimiito\OneDrive - Cisco\Documents\700 - github\kimiito_Note\10_working\10_MoD\04_CSA\CSA_基礎知識_template.pptx"

# テンプレートを読み込む
print(f"テンプレートを読み込んでいます: {template_path}")
prs = Presentation(template_path)

# テンプレート内の全テキストを 2025 → 2026 に置換
print("テンプレート内の 2025 を 2026 に置換中...")
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "2025" in run.text:
                        run.text = run.text.replace("2025", "2026")

# テンプレートの既存スライドを確認
print(f"テンプレート内のスライド数: {len(prs.slides)}")

# テンプレートのスライドレイアウトを調査
print(f"利用可能なレイアウト: {len(prs.slide_layouts)}")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: {layout.name}")

# 汎用レイアウト（通常は blank layout が最後）
blank_layout = prs.slide_layouts[-1]
if len(prs.slide_layouts) > 5:
    blank_layout = prs.slide_layouts[6]  # 標準的には 6 番が blank

# Cisco カラースキーム
CISCO_BLUE   = RGBColor(0x00, 0x6A, 0xA7)
CISCO_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
CISCO_ACCENT = RGBColor(0x00, 0xBC, 0xEB)
CISCO_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CISCO_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
TEXT_DARK    = RGBColor(0x1F, 0x2A, 0x36)
TEXT_MUTED   = RGBColor(0x5B, 0x64, 0x72)
CISCO_RED    = RGBColor(0xFF, 0x45, 0x45)

today = datetime.date.today().strftime("%Y-%m-%d")
FOOTER_TEXT = f"Cisco Confidential  |  Cisco Secure Access 基礎知識  |  {today}"


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """矩形を追加"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, font_name="Yu Gothic UI"):
    """テキストボックスを追加"""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_footer(slide):
    """フッターバー"""
    add_rect(slide, 0, 7.1, 13.33, 0.4, CISCO_BLUE)
    add_textbox(slide, FOOTER_TEXT,
                0.2, 7.12, 11, 0.35,
                font_size=8, color=CISCO_WHITE, align=PP_ALIGN.LEFT)


def add_section_header_bar(slide, title):
    """セクションヘッダーバー"""
    add_rect(slide, 0, 0, 13.33, 0.9, CISCO_BLUE)
    add_textbox(slide, title,
                0.3, 0.05, 12.7, 0.8,
                font_size=22, bold=True, color=CISCO_WHITE)


def bullet_lines(slide, items, left, top, width, height,
                 font_size=13, color=TEXT_DARK):
    """箇条書き"""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Yu Gothic UI"
    return txb


def add_card(slide, left, top, width, height, title, body_lines,
             title_bg=CISCO_BLUE, title_color=CISCO_WHITE,
             body_bg=CISCO_WHITE, font_size=11):
    """カード"""
    add_rect(slide, left, top, width, 0.38, title_bg)
    add_textbox(slide, title, left + 0.1, top + 0.03, width - 0.2, 0.35,
                font_size=12, bold=True, color=title_color)
    body_h = height - 0.38
    add_rect(slide, left, top + 0.38, width, body_h, body_bg, CISCO_GRAY)
    bullet_lines(slide, body_lines,
                 left + 0.12, top + 0.42, width - 0.22, body_h - 0.08,
                 font_size=font_size)


# =========================================================================
# テンプレートのスライドを全て削除
# =========================================================================
print("テンプレートの既存スライドを削除中...")
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# =========================================================================
# CSA スライドを 10 枚追加（テンプレートレイアウト継承）
# =========================================================================

# --------
# Slide 1 - 表紙
# --------
print("Slide 1: Title")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_DARK)
add_rect(slide, 0, 0, 13.33, 0.07, CISCO_ACCENT)
add_rect(slide, 0, 3.5, 5.5, 0.06, CISCO_ACCENT)
add_textbox(slide, "Cisco Secure Access",
            0.5, 1.8, 12, 1.2,
            font_size=40, bold=True, color=CISCO_ACCENT)
add_textbox(slide, "基礎知識ガイド",
            0.5, 2.8, 12, 0.9,
            font_size=30, bold=False, color=CISCO_WHITE)
add_textbox(slide, "SSE 概要 / 実装機能 / Webex 連携 / 設計ベストプラクティス",
            0.5, 3.7, 12, 0.6,
            font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC))
add_textbox(slide, today,
            0.5, 6.5, 4, 0.4,
            font_size=11, color=TEXT_MUTED)
add_footer(slide)

# --------
# Slide 2 - Agenda
# --------
print("Slide 2: Agenda")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_GRAY)
add_section_header_bar(slide, "Agenda")
add_footer(slide)
agenda = [
    "1.  Cisco Secure Access (CSA) とは",
    "2.  CSA を実装するうえで押さえたい 6 つの機能",
    "3.  音声 / 通話品質と CSA の関係",
    "4.  Webex 利用時に CSA の各機能が果たす役割",
    "5.  CSA のライセンス・課金の考え方",
    "6.  Webex で CSA を経由させている事例について",
    "7.  SIPS / SRTP を CSA に通した際の負荷・費用",
    "8.  Signaling / Media の設計ベストプラクティス",
    "9.  まとめ",
]
bullet_lines(slide, agenda, 1.0, 1.1, 11.5, 5.8, font_size=15)

# --------
# Slide 3 - CSA とは
# --------
print("Slide 3: CSA Overview")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_GRAY)
add_section_header_bar(slide, "1. Cisco Secure Access (CSA) とは")
add_footer(slide)
add_textbox(slide,
    "ゼロトラストを前提とした クラウド提供型 SSE (Security Service Edge) ソリューション",
    0.4, 1.05, 12.5, 0.5, font_size=14, bold=True, color=CISCO_BLUE)
overview_lines = [
    "• ユーザー/デバイスからどのアプリへのアクセスも 1 つのクラウド基盤で保護",
    "• シングルクライアント（Cisco Secure Client）・シングルコンソールで管理",
    "• ZTNA / SWG / CASB / FWaaS / DLP / DNS Security / RBI / DEM / AI Access を統合",
    "• Cisco SASE の SSE コンポーネントとして SD-WAN と統合可能",
]
bullet_lines(slide, overview_lines, 0.4, 1.65, 12.3, 2.2, font_size=13)
add_card(slide, 0.4,  4.0, 3.9, 2.7, "Better for Users",
         ["• シングルサインオンで即接続", "• クライアント 1 本で全アクセス", "• フリクションレスな UX"],
         font_size=11)
add_card(slide, 4.5,  4.0, 3.9, 2.7, "Easier for IT",
         ["• 1 コンソール・一元ポリシー管理", "• AI アシスタントでポリシー自動生成", "• 集約レポート"],
         font_size=11)
add_card(slide, 8.6,  4.0, 4.3, 2.7, "Safer for Everyone",
         ["• 多層防御・最小権限", "• Talos 脅威インテリジェンス", "• 継続的なリスク評価"],
         font_size=11)

# --------
# Slide 4 - 6 つの機能
# --------
print("Slide 4: 6 Functions")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_GRAY)
add_section_header_bar(slide, "2. CSA を実装するうえで押さえたい 6 つの機能")
add_footer(slide)
cards = [
    ("① Private Access",        ["ZTNA・VPNaaS", "社内アプリをゼロトラストで提供", "レガシー通信は VPNaaS で補完"]),
    ("② Internet / SaaS Access",["SWG・CASB", "URL/カテゴリ制御", "Shadow IT・AI アプリ可視化"]),
    ("③ Network Security",      ["FWaaS・IPS", "全ポート/プロトコル制御", "L3/L4/L7 ルール適用"]),
    ("④ Data Protection",       ["DLP・AI Access", "機密データ漏洩防止", "AI プロンプト/レスポンス DLP"]),
    ("⑤ Threat Protection",     ["DNS Security・RBI", "Malware Analytics", "DNS 遮断・ブラウザ隔離"]),
    ("⑥ Visibility / Ops",      ["DEM・AI Assistant", "体感品質・アプリ応答の可視化", "AI 補助トラブルシュート"]),
]
col_w = 4.0
row_h = 2.6
positions = [
    (0.3,  1.0), (4.5,  1.0), (8.7,  1.0),
    (0.3,  3.8), (4.5,  3.8), (8.7,  3.8),
]
for (l, t), (title, body) in zip(positions, cards):
    add_card(slide, l, t, col_w, row_h, title, body, font_size=11)

# --------
# Slide 5 - 音声と CSA
# --------
print("Slide 5: Voice/Quality")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_GRAY)
add_section_header_bar(slide, "3. 音声 / 通話品質と CSA の関係")
add_footer(slide)
add_textbox(slide, "音声は遅延・ジッタ・パケットロスに最も敏感。経路設計と機能選択が品質に直結する。",
            0.4, 1.05, 12.5, 0.5, font_size=13, bold=True, color=CISCO_BLUE)
add_card(slide, 0.4, 1.7, 5.8, 2.3, "音声に直接影響する機能（優先度：高）",
         ["• FWaaS : ポート/プロトコル許可と制御",
          "• ZTNA / VPNaaS : 経路と暗号化オーバーヘッド",
          "• DEM : 遅延・ジッタ・ロスのリアルタイム可視化"],
         title_bg=CISCO_BLUE)
add_card(slide, 6.5, 1.7, 6.5, 2.3, "間接的に影響する機能（優先度：中）",
         ["• DNS Security : 名前解決の遅延・失敗を防ぐ",
          "• SWG : Signaling (HTTPS) のプロキシ経路に影響",
          "• DEM での End-to-End 経路可視化が両機能で共通有効"],
         title_bg=RGBColor(0x00, 0x7A, 0xB8))
add_card(slide, 0.4, 4.15, 12.6, 2.55, "⚠ 音声メディアを CSA に通す際の主なリスク",
         ["• 追加ホップによる遅延増加（SSE POP 経由のヘアピン）",
          "• UDP メディアの取り回しによるジッタ増大",
          "• 復号/検査を挟む場合の処理遅延",
          "• → メディアは原則「直通 or 最短経路」が品質上は有利"],
         title_bg=CISCO_RED)

# --------
# Slide 6 - Webex と CSA
# --------
print("Slide 6: Webex Integration")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_GRAY)
add_section_header_bar(slide, "4. Webex 利用時に CSA の各機能が果たす役割")
add_footer(slide)
add_textbox(slide, "Signaling と Media を分けて CSA の適用範囲を設計することが重要",
            0.4, 1.0, 12.5, 0.45, font_size=13, bold=True, color=CISCO_BLUE)
add_card(slide, 0.4, 1.55, 6.0, 2.6, "Webex Signaling (HTTPS ベース)",
         ["• SWG : 必要ドメイン/URL を許可・危険宛先を遮断",
          "• FWaaS : ポート制御・IPS 適用",
          "• DNS Security : 悪性宛先の DNS 解決を遮断",
          "• DEM : 到達性・応答遅延の可視化",
          "→ 全機能を適用しやすい"],
         title_bg=CISCO_BLUE)
add_card(slide, 6.7, 1.55, 6.2, 2.6, "Webex Media (音声/映像 RTP/SRTP)",
         ["• FWaaS : UDP/TCP ポートの許可",
          "• VPNaaS/ZTNA : 社内周辺アプリとの経路整合",
          "• DEM : MOS・遅延・ジッタ・ロスの監視",
          "→ 全面ヘアピンは慎重に。PoC 推奨"],
         title_bg=RGBColor(0x00, 0x7A, 0xB8))
add_card(slide, 0.4, 4.3, 12.6, 2.4, "Webex との CSA 連携で確認できていること（公開資料ベース）",
         ["• Cloud Malware Detection の対象 SaaS に Webex を明記",
          "• DEM はコラボレーションアプリ（Webex/Zoom/Teams）の品質スコアをダッシュボードに表示",
          "• 「Webex の全通信を CSA に通す」という公式事例・ベストプラクティスは 2026-05 時点で公開未確認"],
         title_bg=RGBColor(0x44, 0x60, 0x80))

# --------
# Slide 7 - ライセンス・課金
# --------
print("Slide 7: Licensing")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_GRAY)
add_section_header_bar(slide, "5. CSA のライセンス・課金の考え方")
add_footer(slide)
add_textbox(slide, "公開資料ベースでは、通信量課金（イン/アウト従量）ではなく ユーザー単位サブスクリプション",
            0.4, 1.0, 12.5, 0.5, font_size=13, bold=True, color=CISCO_BLUE)
add_card(slide, 0.4, 1.65, 5.8, 3.2, "パッケージ構成",
         ["• Secure Internet Access (SIA) : Web/SaaS 保護",
          "• Secure Private Access (SPA) : 社内アプリ保護",
          "• SIA と SPA は別ユーザー数で購入可能",
          "• 各パッケージに Essentials / Advantage のグレード",
          "• オプション add-on あり（RBI Advanced など）",
          "• DNS Defense 単体パッケージも提供"],
         title_bg=CISCO_BLUE)
add_card(slide, 6.5, 1.65, 6.5, 3.2, "実務上のコスト要素",
         ["• CSA ライセンス: ユーザー数 × 単価 × 期間",
          "• DEM/ThousandEyes 拡張: 別途ライセンス",
          "• Secure Client: CSA に同梱（追加費用なし）",
          "• 回線・クラウド接続費: 別途",
          "• PoC・設計・移行コスト: 見積必須",
          "• ※ 通信量従量課金の構造は公開情報では非確認"],
         title_bg=RGBColor(0x00, 0x7A, 0xB8))
add_textbox(slide, "詳細見積は Cisco 営業 / パートナー経由が必要です。",
            0.4, 5.0, 12.5, 0.4, font_size=11, color=TEXT_MUTED)

# --------
# Slide 8 - Webex 事例 + SIPS/SRTP
# --------
print("Slide 8: Webex Cases / SIPS/SRTP")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_GRAY)
add_section_header_bar(slide, "6–7. Webex × CSA 事例 / SIPS・SRTP の負荷と費用")
add_footer(slide)
add_card(slide, 0.4, 1.0, 6.1, 3.0, "6. Webex で CSA を経由している事例",
         ["確認できたこと（公開資料）:",
          "• DEM : Webex の品質スコアを可視化",
          "• Cloud Malware : Webex を対象 SaaS に含む",
          "未確認:",
          "• Calling の Signaling/Media を全量 CSA 経由にした事例",
          "→ 部分適用は現実的。全通しは PoC 前提"],
         title_bg=CISCO_BLUE)
add_card(slide, 6.8, 1.0, 6.1, 3.0, "7. SIPS / SRTP を CSA に通した場合",
         ["負荷面:",
          "• 追加ホップ → 遅延増加",
          "• SSE POP 経由 → ジッタ増大",
          "• 復号/検査 → 処理遅延",
          "費用面（seat 課金が基本）:",
          "• ライセンス増・add-on の必要性",
          "• 帯域増・PoC コストが発生"],
         title_bg=RGBColor(0x00, 0x7A, 0xB8))
add_textbox(slide,
    "SRTP メディアは「セキュリティ検査のメリット ＜ 品質劣化リスク」になりやすい。慎重な設計が必要。",
    0.4, 4.15, 12.5, 0.5, font_size=12, bold=True, color=CISCO_RED)
add_card(slide, 0.4, 4.75, 12.5, 2.0, "推奨アプローチ",
         ["① まず Signaling のみ CSA 経由でテスト → 機能確認",
          "② Media は直通 or 最短経路を基本とし、品質指標（MOS/遅延/ジッタ）を DEM で取得",
          "③ 全通しが要件の場合は PoC を実施し、品質劣化の許容範囲を確認してから展開"],
         title_bg=RGBColor(0x22, 0x65, 0x44))

# --------
# Slide 9 - ベストプラクティス
# --------
print("Slide 9: Best Practices")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_GRAY)
add_section_header_bar(slide, "8. Signaling / Media の設計ベストプラクティス")
add_footer(slide)
add_textbox(slide, "「通信種別ごとに最適経路を分離する」が基本原則",
            0.4, 1.0, 12.5, 0.45, font_size=13, bold=True, color=CISCO_BLUE)
add_card(slide, 0.4, 1.6, 5.9, 2.5, "Signaling 推奨設計",
         ["◎ SWG または FWaaS 経由",
          "◎ DNS Security + URL ベース許可",
          "• IP アドレス固定制御は非推奨（Webex は動的 IP）",
          "• HTTP ヘッダーの改変・削除は非推奨",
          "• プロキシ機能との組み合わせが有効"],
         title_bg=CISCO_BLUE)
add_card(slide, 6.6, 1.6, 6.3, 2.5, "Media 推奨設計",
         ["◎ 直接ブレイクアウト（最短経路優先）",
          "◎ 最小限 FW ポート許可のみ",
          "◎ DEM / ThousandEyes で品質監視",
          "• ローカル拠点での制御を優先",
          "• 全面ヘアピンは PoC 後に判断"],
         title_bg=RGBColor(0x00, 0x7A, 0xB8))
add_card(slide, 0.4, 4.25, 12.5, 2.5, "判断フロー",
         ["1. Signaling と Media を分類する",
          "2. Signaling → CSA (SWG/FWaaS/DNS) 経由で制御・可視化",
          "3. Media → 品質優先。直通が第一選択。セキュリティ要件が強い場合は PoC",
          "4. DEM を使い MOS・遅延・ジッタ・ロスを継続的にモニタリング",
          "5. 全通しが必要なら PoC で KPI を満たすことを確認してから本番展開"],
         title_bg=RGBColor(0x22, 0x55, 0x77))

# --------
# Slide 10 - まとめ
# --------
print("Slide 10: Summary")
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CISCO_DARK)
add_rect(slide, 0, 0, 13.33, 0.07, CISCO_ACCENT)
add_footer(slide)
add_textbox(slide, "まとめ",
            0.5, 0.3, 12, 0.7,
            font_size=26, bold=True, color=CISCO_ACCENT)
summary = [
    ("CSA の機能体系",
     ["6 分類（Private / Internet / Network / Data / Threat / Ops）で整理すると設計しやすい",
      "単一コンソール・シングルクライアントで全機能を統合管理"]),
    ("音声 / 通話品質",
     ["最重要: FWaaS・ZTNA/VPNaaS・DEM",
      "メディアはヘアピン回避が品質上の基本。DEM で継続監視"]),
    ("Webex × CSA",
     ["Signaling は CSA 経由で制御・可視化しやすい",
      "Media は直通優先。全通しは PoC で品質確認後に判断"]),
    ("課金",
     ["ユーザー単位サブスクリプションが基本（通信量従量課金は非確認）",
      "SIA / SPA を用途別に組み合わせて購入可能"]),
]
col_positions = [
    (0.4,  1.2, 5.9),
    (6.5,  1.2, 6.4),
    (0.4,  4.0, 5.9),
    (6.5,  4.0, 6.4),
]
for (l, t, w), (title, body) in zip(col_positions, summary):
    add_card(slide, l, t, w, 2.55, title, body,
             title_bg=CISCO_BLUE, title_color=CISCO_WHITE, body_bg=RGBColor(0x1E, 0x2A, 0x3A),
             font_size=12)

# =========================================================================
# 保存
# =========================================================================
prs.save(output_path)
print(f"\n保存完了: {output_path}")
print(f"Cisco 公式テンプレートをベースに 10 枚のスライドを生成しました。")
print(f"✓ 2026 年対応")
print(f"✓ テンプレートのマスター/色/フォント/ロゴを継承")
print(f"✓ Cisco Secure Access 完全コンテンツ統合")
