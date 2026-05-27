from pathlib import Path
import datetime
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

base_dir = Path(r"c:\Users\kimiito\OneDrive - Cisco\Documents\700 - github\kimiito_Note\10_working\10_MoD\04_CSA")
template_path = base_dir / "cisco2025_template.pptx"
out_path = base_dir / "CSA_基礎知識_template_native.pptx"

prs = Presentation(str(template_path))


def replace_year_in_shape(shape):
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            if "2025" in r.text:
                r.text = r.text.replace("2025", "2026")


def replace_year_everywhere(p):
    for s in p.slides:
        for sh in s.shapes:
            replace_year_in_shape(sh)

    for master in p.slide_masters:
        for sh in master.shapes:
            replace_year_in_shape(sh)
        for layout in master.slide_layouts:
            for sh in layout.shapes:
                replace_year_in_shape(sh)


def get_body_placeholders(slide):
    bodies = []
    for ph in slide.placeholders:
        if ph.is_placeholder and ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
            bodies.append(ph)
    bodies.sort(key=lambda x: x.placeholder_format.idx)
    return bodies


def set_title(slide, text):
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
        return
    for ph in slide.placeholders:
        if ph.is_placeholder and ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            ph.text = text
            return


def set_bullets(ph, lines):
    tf = ph.text_frame
    tf.clear()
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0


def set_subtitle(slide, lines):
    bodies = get_body_placeholders(slide)
    if bodies:
        set_bullets(bodies[0], lines)


# Cisco light 5 12 2025 は 2つ目のマスターに含まれる想定
if len(prs.slide_masters) < 2:
    raise RuntimeError("期待した Cisco light マスターが見つかりませんでした")

light_master = prs.slide_masters[1]
layouts = light_master.slide_layouts

# Layout mapping in Cisco light master
L_TITLE = layouts[0]      # Title Slide 1, Two Speakers
L_AGENDA = layouts[8]     # Agenda 1
L_TWO_COL = layouts[19]   # Title, Subtitle, 2 Columns
L_THREE_COL = layouts[20] # Title, Subtitle, 3 Columns
L_BULLET = layouts[54]    # Bullet
L_THANKYOU = layouts[51]  # Thank you 1

replace_year_everywhere(prs)

# Remove existing sample slides but keep masters/themes
while len(prs.slides) > 0:
    rid = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rid)
    del prs.slides._sldIdLst[0]

# Slide 1: Title
slide = prs.slides.add_slide(L_TITLE)
set_title(slide, "Cisco Secure Access 基礎知識ガイド")
set_subtitle(slide, [
    "SSE 概要 / 実装機能 / Webex 連携 / 設計ベストプラクティス",
    datetime.date.today().strftime("%Y-%m-%d"),
    "Cisco light 5 12 2026",
])

# Slide 2: Agenda
slide = prs.slides.add_slide(L_AGENDA)
set_title(slide, "Agenda")
bodies = get_body_placeholders(slide)
if len(bodies) >= 1:
    set_bullets(bodies[0], [
        "1. Cisco Secure Access (CSA) とは",
        "2. CSA を実装するうえで押さえたい 6 つの機能",
        "3. 音声 / 通話品質と CSA の関係",
        "4. Webex 利用時に CSA の各機能が果たす役割",
        "5. CSA のライセンス・課金の考え方",
        "6. Webex で CSA を経由させている事例について",
        "7. SIPS / SRTP を CSA に通した際の負荷・費用",
        "8. Signaling / Media の設計ベストプラクティス",
        "9. まとめ",
    ])

# Slide 3: Overview (2 columns)
slide = prs.slides.add_slide(L_TWO_COL)
set_title(slide, "1. Cisco Secure Access (CSA) とは")
bodies = get_body_placeholders(slide)
if len(bodies) >= 3:
    set_bullets(bodies[0], ["ゼロトラスト前提のクラウド提供型 SSE ソリューション"])
    set_bullets(bodies[1], [
        "ユーザー/デバイスからアプリへのアクセスを統合保護",
        "シングルクライアント・シングルコンソールで運用",
        "Cisco SASE の SSE コンポーネントとして SD-WAN と統合可能",
    ])
    set_bullets(bodies[2], [
        "統合機能:",
        "ZTNA / SWG / CASB / FWaaS",
        "DLP / DNS Security / RBI / DEM / AI Access",
        "期待効果:",
        "Better for Users / Easier for IT / Safer for Everyone",
    ])

# Slide 4: 6 functions (3 columns)
slide = prs.slides.add_slide(L_THREE_COL)
set_title(slide, "2. CSA を実装するうえで押さえたい 6 つの機能")
bodies = get_body_placeholders(slide)
if len(bodies) >= 4:
    set_bullets(bodies[0], ["6機能を段階導入"])
    set_bullets(bodies[1], [
        "1) Private Access (ZTNA/VPNaaS)",
        "2) Internet/SaaS Access (SWG/CASB)",
    ])
    set_bullets(bodies[2], [
        "3) Network Security (FWaaS/IPS)",
        "4) Data Protection (DLP/AI Access)",
    ])
    set_bullets(bodies[3], [
        "5) Threat Protection (DNS Security/RBI)",
        "6) Visibility/Ops (DEM/AI Assistant)",
    ])

# Slide 5: Voice quality
slide = prs.slides.add_slide(L_TWO_COL)
set_title(slide, "3. 音声 / 通話品質と CSA の関係")
bodies = get_body_placeholders(slide)
if len(bodies) >= 3:
    set_bullets(bodies[0], ["音声は遅延・ジッタ・ロスに最も敏感"])
    set_bullets(bodies[1], [
        "直接影響(優先度高): FWaaS / ZTNA・VPNaaS / DEM",
        "間接影響(優先度中): DNS Security / SWG",
    ])
    set_bullets(bodies[2], [
        "メディア経路は直通 or 最短経路を原則",
        "全面ヘアピンは PoC で品質評価後に判断",
    ])

# Slide 6: Webex integration
slide = prs.slides.add_slide(L_TWO_COL)
set_title(slide, "4. Webex 利用時に CSA の各機能が果たす役割")
bodies = get_body_placeholders(slide)
if len(bodies) >= 3:
    set_bullets(bodies[0], ["Signaling と Media を分離設計"])
    set_bullets(bodies[1], [
        "Signaling(HTTPS): SWG/FWaaS/DNS/DEM を適用しやすい",
        "制御・可視化の中心",
    ])
    set_bullets(bodies[2], [
        "Media(RTP/SRTP): FW最小許可 + DEM監視",
        "全面経由は PoC を前提に検証",
    ])

# Slide 7: Licensing
slide = prs.slides.add_slide(L_TWO_COL)
set_title(slide, "5. CSA のライセンス・課金の考え方")
bodies = get_body_placeholders(slide)
if len(bodies) >= 3:
    set_bullets(bodies[0], ["基本はユーザー単位サブスクリプション"])
    set_bullets(bodies[1], [
        "SIA / SPA は別ユーザー数で購入可能",
        "Essentials / Advantage グレード",
    ])
    set_bullets(bodies[2], [
        "コスト: ライセンス + 拡張 + 回線 + PoC/移行",
        "通信量従量課金モデルは公開情報で未確認",
    ])

# Slide 8: Case and SIPS/SRTP
slide = prs.slides.add_slide(L_TWO_COL)
set_title(slide, "6-7. Webex 事例 / SIPS・SRTP の負荷と費用")
bodies = get_body_placeholders(slide)
if len(bodies) >= 3:
    set_bullets(bodies[0], ["公開情報ベースの確認事項"])
    set_bullets(bodies[1], [
        "DEM で Webex 品質可視化",
        "Cloud Malware Detection 対象に Webex 記載",
        "Calling 全量経由の公式事例は未確認",
    ])
    set_bullets(bodies[2], [
        "SIPS/SRTP は遅延・ジッタ増のリスク",
        "推奨: Signaling 先行、Media は PoC で判断",
    ])

# Slide 9: Best practices
slide = prs.slides.add_slide(L_TWO_COL)
set_title(slide, "8. Signaling / Media の設計ベストプラクティス")
bodies = get_body_placeholders(slide)
if len(bodies) >= 3:
    set_bullets(bodies[0], ["通信種別ごとの経路最適化"])
    set_bullets(bodies[1], [
        "Signaling: SWG/FWaaS + DNS Security",
        "URLベース許可、ヘッダー改変は避ける",
    ])
    set_bullets(bodies[2], [
        "Media: 直接ブレイクアウト + DEM監視",
        "全通しは KPI 定義のうえ PoC で検証",
    ])

# Slide 10: Summary
slide = prs.slides.add_slide(L_BULLET)
set_title(slide, "まとめ")
bodies = get_body_placeholders(slide)
if bodies:
    set_bullets(bodies[0], [
        "CSA は 6 機能統合でゼロトラストアクセスを実現",
        "Webex は Signaling/Media を分離して設計",
        "Media は品質優先、全通しは PoC 判断",
        "課金はユーザー単位サブスクリプションが基本",
        "本資料は 2026 年版に更新済み",
    ])

replace_year_everywhere(prs)
prs.save(str(out_path))

print(f"保存完了: {out_path}")
print("Cisco light 5 12 2025 マスター(2つ目)で 10 枚を再生成しました")
