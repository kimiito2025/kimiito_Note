from pptx import Presentation
p = Presentation(r"c:\Users\kimiito\OneDrive - Cisco\Documents\700 - github\kimiito_Note\10_working\10_MoD\04_CSA\cisco2025_template.pptx")
print('slides', len(p.slides))
for i, l in enumerate(p.slide_layouts):
    print(f"LAYOUT {i}: {l.name}")
    for ph in l.placeholders:
        try:
            print('  ph', ph.placeholder_format.idx, ph.name, ph.placeholder_format.type)
        except Exception:
            print('  ph ?', ph.name)
