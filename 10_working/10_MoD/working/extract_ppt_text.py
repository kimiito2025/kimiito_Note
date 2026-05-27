from pptx import Presentation
from pathlib import Path
import json

ppt = Path(r"c:\Users\kimiito\OneDrive - Cisco\Documents\700 - github\kimiito_Note\10_working\10_MoD\04_CSA\CSA_基礎知識.pptx")
prs = Presentation(str(ppt))
out = []
for i, s in enumerate(prs.slides, start=1):
    title = ""
    if s.shapes.title is not None and hasattr(s.shapes.title, 'text'):
        title = s.shapes.title.text.strip()
    texts = []
    for sh in s.shapes:
        if hasattr(sh, "text") and sh.text:
            t = sh.text.strip()
            if t:
                texts.append(t)
    out.append({"slide": i, "title": title, "texts": texts})

print(json.dumps(out, ensure_ascii=False, indent=2))
