from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
p = Presentation(r"c:\Users\kimiito\OneDrive - Cisco\Documents\700 - github\kimiito_Note\10_working\10_MoD\04_CSA\cisco2025_template.pptx")
m = p.slide_masters[1]
print('master1 layouts', len(m.slide_layouts))
for i in [0,8,9,10,11,12,13,17,18,19,20,21,22,41,47,51,53,54]:
    l = m.slide_layouts[i]
    print(f"\nLAYOUT {i}: {l.name}")
    for ph in l.placeholders:
        try:
            print('  idx', ph.placeholder_format.idx, 'type', ph.placeholder_format.type, 'name', ph.name)
        except Exception:
            print('  idx ? name', ph.name)
